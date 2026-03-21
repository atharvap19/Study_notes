import os
import re
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook

ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx"}


def extract_text(file_path: str) -> str:
    """Extract text content from a supported file type (flat string, for backward compat)."""
    chunks = extract_chunks(file_path)
    return format_chunks(chunks)


def extract_chunks(file_path: str) -> list[dict]:
    """
    Extract text as structured chunks with metadata.
    Each chunk: {"section": str, "source": str, "content": str}
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _chunk_pdf(file_path)
    elif ext == ".pptx":
        return _chunk_pptx(file_path)
    elif ext == ".docx":
        return _chunk_docx(file_path)
    elif ext == ".xlsx":
        return _chunk_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def format_chunks(chunks: list[dict]) -> str:
    """Format chunks into structured text with section markers for the LLM."""
    parts = []
    for chunk in chunks:
        header = f"[{chunk['source']}] {chunk['section']}"
        parts.append(f"--- {header} ---\n{chunk['content']}")
    return "\n\n".join(parts)


# --- PDF: chunk by detected headings or by page ---

def _chunk_pdf(path: str) -> list[dict]:
    reader = PdfReader(path)
    chunks = []

    for page_num, page in enumerate(reader.pages, 1):
        page_text = page.extract_text()
        if not page_text or not page_text.strip():
            continue

        # Try to split page by headings (lines that look like section titles)
        sections = _split_by_headings(page_text)

        if len(sections) > 1:
            # Multiple sections found on this page
            for title, content in sections:
                if content.strip():
                    chunks.append({
                        "section": title,
                        "source": f"Page {page_num}",
                        "content": content.strip(),
                    })
        else:
            # No headings detected — use whole page as one chunk
            chunks.append({
                "section": f"Page {page_num} Content",
                "source": f"Page {page_num}",
                "content": page_text.strip(),
            })

    return chunks


def _split_by_headings(text: str) -> list[tuple[str, str]]:
    """
    Detect heading-like lines and split text into (heading, content) pairs.
    Headings are detected as:
    - Numbered sections (1. Introduction, 3.1 Data Preprocessing)
    - ALL CAPS lines
    - Short lines (< 80 chars) followed by longer content
    """
    lines = text.split("\n")
    sections = []
    current_title = "Untitled Section"
    current_content = []

    heading_pattern = re.compile(
        r"^(\d+\.?\d*\.?\s+.{3,80})$"  # Numbered: "1. Introduction", "3.1 Data"
    )

    for line in lines:
        stripped = line.strip()
        if not stripped:
            current_content.append("")
            continue

        is_heading = False

        # Check numbered heading pattern
        if heading_pattern.match(stripped):
            is_heading = True
        # Check ALL CAPS heading (at least 3 words)
        elif stripped.isupper() and len(stripped.split()) >= 2 and len(stripped) < 80:
            is_heading = True

        if is_heading:
            # Save previous section
            if current_content:
                sections.append((current_title, "\n".join(current_content)))
            current_title = stripped
            current_content = []
        else:
            current_content.append(stripped)

    # Don't forget last section
    if current_content:
        sections.append((current_title, "\n".join(current_content)))

    return sections


# --- PPTX: each slide is a chunk, with slide title as section ---

def _chunk_pptx(path: str) -> list[dict]:
    prs = Presentation(path)
    chunks = []

    for slide_num, slide in enumerate(prs.slides, 1):
        title = None
        body_text = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if not line:
                        continue
                    # First text found on slide with title placeholder = title
                    if title is None and shape.shape_type is not None:
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                            if shape.placeholder_format.idx == 0:  # Title placeholder
                                title = line
                                continue
                    body_text.append(line)

        # If no title placeholder found, use first line as title
        if title is None and body_text:
            title = body_text.pop(0)

        content = "\n".join(body_text) if body_text else ""
        if title or content:
            chunks.append({
                "section": title or f"Slide {slide_num}",
                "source": f"Slide {slide_num}",
                "content": content if content else title or "",
            })

    return chunks


# --- DOCX: chunk by heading styles ---

def _chunk_docx(path: str) -> list[dict]:
    doc = Document(path)
    chunks = []
    current_title = "Document Start"
    current_content = []
    section_count = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect heading styles (Heading 1, Heading 2, etc.)
        is_heading = para.style.name.startswith("Heading")

        if is_heading:
            # Save previous section
            if current_content:
                section_count += 1
                chunks.append({
                    "section": current_title,
                    "source": f"Section {section_count}",
                    "content": "\n".join(current_content),
                })
            current_title = text
            current_content = []
        else:
            current_content.append(text)

    # Last section
    if current_content:
        section_count += 1
        chunks.append({
            "section": current_title,
            "source": f"Section {section_count}",
            "content": "\n".join(current_content),
        })

    return chunks


# --- XLSX: each sheet is a chunk ---

def _chunk_xlsx(path: str) -> list[dict]:
    wb = load_workbook(path, data_only=True)
    chunks = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames, 1):
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))

        if rows:
            chunks.append({
                "section": sheet_name,
                "source": f"Sheet {sheet_idx}",
                "content": "\n".join(rows),
            })

    return chunks
